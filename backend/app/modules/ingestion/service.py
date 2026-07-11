import asyncio
import logging
import os
import tempfile
import time

from pypdf import PdfReader
from pypdf.errors import PdfReadError

logger = logging.getLogger(__name__)

MAX_PAGES = 50
_docling_converters = {}
_hybrid_chunker = None

# Minimum chars of extractable text for a PDF to count as "has a text layer".
# A scanned/photographed PDF yields ~0; a real text PDF yields hundreds+.
_TEXT_LAYER_MIN_CHARS = 20


def _ocr_blocked() -> bool:
    """When BLOCK_OCR_DOCUMENTS is on, scanned PDFs (no text layer) are rejected
    up front and the docling pipeline runs WITHOUT OCR — so no memory-heavy
    full-page rasterization. Used on constrained hosts (Railway trial) where
    full-page Hebrew OCR OOM-kills the container."""
    return os.getenv("BLOCK_OCR_DOCUMENTS", "false").strip().lower() in ("1", "true", "yes", "on")


def _pdf_has_text_layer(reader) -> bool:
    """True if the PDF has an extractable text layer (i.e. does NOT need OCR).
    Cheap pypdf check — no rendering, no torch — so it can gate the expensive
    docling path before it ever loads."""
    total = 0
    for page in reader.pages:
        try:
            total += len((page.extract_text() or "").strip())
        except Exception:
            continue
        if total >= _TEXT_LAYER_MIN_CHARS:
            return True
    return False


def _lightweight_chunks(text: str, max_chars: int = 1500) -> list[str]:
    """Group plain text into ~max_chars contiguous chunks on sentence/paragraph
    boundaries — a cheap stand-in for HybridChunker when docling is skipped.

    The analyzer makes ONE LLM call per chunk, so without chunking a big doc
    becomes hundreds of per-sentence calls that blow the 180s cap. Chunks are
    exact, in-order substrings of `text`, so hybrid_detector._locate_chunks
    finds them by exact match."""
    import re
    if not text.strip():
        return []
    # Cut points: sentence-ending punctuation (Latin + Hebrew maqaf/sof pasuq)
    # followed by space, or a paragraph break.
    cuts = [m.end() for m in re.finditer(r'(?:[.!?׃׀…]["\')\]]?\s+|\n+)', text)]
    if not cuts or cuts[-1] != len(text):
        cuts.append(len(text))
    chunks = []
    start = 0
    prev = 0
    for c in cuts:
        if c - start > max_chars and prev > start:
            chunks.append(text[start:prev])
            start = prev
        prev = c
    if start < len(text):
        chunks.append(text[start:])
    return [c for c in chunks if c.strip()]


def _get_hybrid_chunker():
    global _hybrid_chunker
    if _hybrid_chunker is None:
        from docling.chunking import HybridChunker
        _hybrid_chunker = HybridChunker()
        logger.info("Docling HybridChunker initialized")
    return _hybrid_chunker

def _get_docling_converter(use_ocr: bool = True):
    if use_ocr not in _docling_converters:
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling.datamodel.pipeline_options import PdfPipelineOptions, TesseractCliOcrOptions
        from docling.datamodel.base_models import InputFormat

        pipeline_opts = PdfPipelineOptions()
        # OCR scanned pages with Tesseract (Hebrew + English). Docling's default
        # engine (EasyOCR) has no Hebrew model, so it must not be used here.
        pipeline_opts.do_ocr = use_ocr
        if use_ocr:
            # force_full_page_ocr: scanned pages are one big image — without it,
            # only small detected bitmap fragments get OCR'd (~0 text per page).
            # This is also the memory-heavy path (rasterizes every page).
            pipeline_opts.ocr_options = TesseractCliOcrOptions(
                lang=["heb", "eng"], force_full_page_ocr=True
            )
        pipeline_opts.do_table_structure = True

        _docling_converters[use_ocr] = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_opts),
            }
        )
        logger.info("Docling converter initialized (ocr=%s)", use_ocr)
    return _docling_converters[use_ocr]

def _extract_docx_fallback(file_bytes: bytes) -> dict:
    """Extract text from DOCX using python-docx when docling is unavailable."""
    import io
    from docx import Document as DocxDocument
    doc = DocxDocument(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    full_text = "\n".join(paragraphs)
    core = doc.core_properties
    title = core.title or None
    author = core.author or None
    has_hebrew = any('֐' <= c <= '׿' for c in full_text[:500])
    detected = "he" if has_hebrew else "en"
    logger.info("DOCX fallback extraction: paragraphs=%d text_length=%d", len(paragraphs), len(full_text))
    return {"text": full_text, "page_count": 1, "title": title, "author": author, "detected_language": detected}


def _extract_pptx_fallback(file_bytes: bytes) -> dict:
    """Extract text from PPTX using python-pptx when docling is unavailable."""
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    full_text = "\n".join(lines)
    has_hebrew = any('֐' <= c <= '׿' for c in full_text[:500])
    detected = "he" if has_hebrew else "en"
    logger.info("PPTX fallback extraction: slides=%d text_length=%d", len(prs.slides), len(full_text))
    return {"text": full_text, "page_count": len(prs.slides), "title": None, "author": None, "detected_language": detected}


def _extract_pdf_fallback(temp_path: str, file_bytes: bytes, max_pages: int) -> dict:
    """Extract text from PDF using pypdf when docling is unavailable."""
    try:
        reader = PdfReader(temp_path)
        page_count = len(reader.pages)
        if page_count > max_pages:
            return {"error": f"Document exceeds {max_pages} page limit ({page_count} pages)"}
        title = None
        author = None
        if reader.metadata:
            title = reader.metadata.title or None
            author = reader.metadata.author or None
        text_parts = []
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
        full_text = "\n".join(text_parts)
        has_hebrew = any('֐' <= c <= '׿' for c in full_text[:500])
        detected = "he" if has_hebrew else "en"
        logger.info("PDF fallback extraction: pages=%d text_length=%d", page_count, len(full_text))
        return {"text": full_text, "page_count": page_count, "title": title, "author": author, "detected_language": detected}
    except PdfReadError as e:
        msg = str(e).lower()
        if "password" in msg or "encrypted" in msg:
            return {"error": "PDF is password-protected"}
        return {"error": "PDF appears corrupted"}


def _parse_document_sync(file_bytes: bytes, filename: str, max_pages: int = MAX_PAGES) -> dict:
    ext = os.path.splitext(filename)[1].lower() or ".pdf"
    temp_fd, temp_path = tempfile.mkstemp(suffix=ext)
    logger.info("Extraction started: filename=%s ext=%s size_bytes=%d", filename, ext, len(file_bytes))

    try:
        with os.fdopen(temp_fd, 'wb') as f:
            f.write(file_bytes)


        page_count = 0
        pdf_title = None
        pdf_author = None

        # TXT passthrough — Docling adds no value for plain text and often
        # introduces spurious layout artifacts, so decode the bytes directly.
        if ext == ".txt":
            try:
                text = file_bytes.decode("utf-8")
            except UnicodeDecodeError:
                try:
                    text = file_bytes.decode("utf-8-sig")
                except UnicodeDecodeError:
                    text = file_bytes.decode("latin-1", errors="replace")
            has_hebrew = any('\u0590' <= c <= '\u05FF' for c in text[:500])
            detected = "he" if has_hebrew else "en"
            logger.info("Extraction completed (TXT passthrough): filename=%s detected_language=%s text_length=%d", filename, detected, len(text))
            return {
                "text": text,
                "page_count": 1,
                "title": None,
                "author": None,
                "detected_language": detected,
            }

        if ext == ".pdf":
            try:
                reader = PdfReader(temp_path)
                page_count = len(reader.pages)
                logger.info("PDF page check: filename=%s pages=%d max_pages=%d", filename, page_count, max_pages)
                if page_count > max_pages:
                    logger.warning("PDF rejected: page limit exceeded filename=%s pages=%d limit=%d", filename, page_count, max_pages)
                    return {"error": f"Document exceeds {max_pages} page limit ({page_count} pages)"}

                if reader.metadata:
                    if reader.metadata.title:
                        pdf_title = reader.metadata.title
                    if reader.metadata.author:
                        pdf_author = reader.metadata.author
            except PdfReadError as e:
                msg = str(e).lower()
                if "password" in msg or "encrypted" in msg:
                    logger.warning("PDF rejected: password-protected filename=%s", filename)
                    return {"error": "PDF is password-protected"}
                logger.error("PDF corrupted: filename=%s", filename)
                return {"error": "PDF appears corrupted"}
            except Exception:
                logger.error("PDF corrupted: filename=%s", filename)
                return {"error": "PDF appears corrupted"}

            # Protection switch: on constrained hosts, reject scanned PDFs (no
            # text layer → would need memory-heavy OCR) before docling loads.
            if _ocr_blocked() and not _pdf_has_text_layer(reader):
                logger.warning("OCR-required PDF rejected (BLOCK_OCR_DOCUMENTS on): filename=%s", filename)
                return {"error": "This document appears to be scanned and has no extractable text layer (no OCR support)."}


        # Try docling first; fall back to lightweight parsers if not installed
        try:
            import docling  # noqa: F401
            _docling_available = True
        except ImportError:
            _docling_available = False

        # Skip docling's memory-heavy vision pipeline entirely when the guard is
        # on (constrained host) — scanned PDFs are already rejected above, so the
        # remaining docs have a text layer that lightweight parsers read cheaply.
        if not _docling_available or _ocr_blocked():
            reason = "not installed" if not _docling_available else "guard on (BLOCK_OCR_DOCUMENTS)"
            logger.info("Using lightweight parser for %s — docling %s", ext, reason)
            if ext == ".docx":
                result = _extract_docx_fallback(file_bytes)
            elif ext == ".pptx":
                result = _extract_pptx_fallback(file_bytes)
            elif ext == ".pdf":
                result = _extract_pdf_fallback(temp_path, file_bytes, max_pages)
            else:
                return {"error": f"Unsupported format {ext} (docling unavailable)"}
            # Chunk the extracted text so analysis makes one LLM call per chunk
            # (not per sentence) — HybridChunker needs docling, so use the cheap
            # text chunker here.
            if "error" not in result and result.get("text"):
                result["chunks"] = _lightweight_chunks(result["text"])
            return result

        _t0 = time.monotonic()
        logger.info("Docling conversion started: filename=%s ext=%s", filename, ext)
        # With the guard on we've already confirmed a text layer, so run docling
        # WITHOUT OCR — skips full-page rasterization, the OOM source.
        converter = _get_docling_converter(use_ocr=not _ocr_blocked())
        result = converter.convert(temp_path)
        logger.info("Docling conversion completed: filename=%s elapsed_s=%.3f", filename, time.monotonic() - _t0)
        doc_dict = result.document.export_to_dict()
        items = doc_dict.get("texts", [])


        meta = getattr(result.document, 'metadata', None)
        title = getattr(meta, 'title', None) or pdf_title
        author = getattr(meta, 'author', None) or pdf_author


        meaningful_items = [item for item in items if item.get('label') not in ['page_header', 'page_footer']]
        first_few_items = meaningful_items[:20]

        title_index = -1
        if not title:

            for i, item in enumerate(first_few_items):
                if item.get('label') == 'title':
                    title = item.get('text', '')
                    title_index = i
                    break
            

            if not title:
                noise = ["university", "college", "school of", "digitalcommons", "journal", "abstract", "published", "אוניברסיטת", "אוניברסיטה", "מכללת", "מכללה", "בית ספר", "פקולטה", "מחלקה", "תקציר", "פורסם"]
                for i, item in enumerate(first_few_items):
                    text_val = item.get('text', '')
                    text_lower = text_val.lower()

                    if not any(n in text_lower for n in noise) and 10 < len(text_val) < 200:
                        title = text_val.strip()
                        title_index = i
                        break

        if title and title_index == -1:
            for i, item in enumerate(first_few_items):

                if title[:20].lower() in item.get('text', '').lower():
                    title_index = i
                    break

        if not author:

            for item in first_few_items:
                text_val = item.get('text', '').strip()
                lower_text = text_val.lower()
                prefixes = {"by ": 3, "מאת ": 4, "מאת:": 4, "על ידי ": 7, 'נכתב ע"י ': 9, "נכתב על ידי ": 12}
                found_author = False
                for prefix, length in prefixes.items():
                    if lower_text.startswith(prefix):
                        author = text_val[length:].strip()
                        found_author = True
                        break
                if found_author:
                    break
            

            if not author and title_index != -1:
                for i in range(title_index + 1, min(title_index + 4, len(first_few_items))):
                    next_item = first_few_items[i]
                    text_val = next_item.get('text', '').strip()
                    text_lower = text_val.lower()
                    

                    if any(n in text_lower for n in ["abstract", "introduction", "university", "department", "email", "תקציר", "מבוא", "אוניברסיטה", "מחלקה", "דוא\"ל", "אימייל", "הקדמה", "פקולטה"]):
                        continue
                        

                    if 3 < len(text_val) < 80:
                        author = text_val
                        break


        # Use plain text for analysis — HybridChunker chunks are plain text substrings
        # of export_to_text(), not markdown, so this eliminates chunk-location failures.
        full_text = result.document.export_to_text()
        has_hebrew = any('\u0590' <= c <= '\u05FF' for c in full_text[:500])
        detected_lang = "he" if has_hebrew else "en"
        final_pages = len(result.document.pages) if hasattr(result.document, 'pages') else page_count

        chunker = _get_hybrid_chunker()
        chunks_data = [c.text for c in chunker.chunk(result.document) if c.text.strip()]

        # PDF: build offset→bbox index for visual overlay rendering.
        bbox_annotations = None
        page_sizes = None
        if ext == ".pdf":
            bbox_annotations = []
            # Sort items into reading order (page asc, top-of-page first) so the
            # sequential find() search advances monotonically and doesn't miss items
            # when Docling returns them out of order.
            def _reading_order(item):
                prov_list = item.get("prov", [])
                if not prov_list:
                    return (9999, -9999.0)
                prov = prov_list[0]
                pg = prov.get("page_no", 9999)
                raw = prov.get("bbox")
                t = (raw.t if hasattr(raw, 't') else raw.get('t', 0.0)) if raw else 0.0
                return (pg, -t)  # higher t = top of page = earlier in reading order
            search_start = 0
            for item in sorted(items, key=_reading_order):
                item_text = item.get("text", "")
                prov_list = item.get("prov", [])
                if not item_text or not prov_list:
                    continue
                prov = prov_list[0]
                raw_bbox = prov.get("bbox")
                if raw_bbox is None:
                    continue
                bbox = {"l": raw_bbox.l, "t": raw_bbox.t, "r": raw_bbox.r, "b": raw_bbox.b} if hasattr(raw_bbox, 'l') else raw_bbox
                page_no = prov.get("page_no", 1)
                idx = full_text.find(item_text, search_start)
                if idx != -1:
                    bbox_annotations.append({
                        "start": idx,
                        "end": idx + len(item_text),
                        "page": page_no,
                        "bbox": bbox,
                    })
                    search_start = idx + len(item_text)

            page_sizes = {}
            if hasattr(result.document, 'pages'):
                for pg_no, pg in result.document.pages.items():
                    sz = getattr(pg, 'size', None)
                    if sz:
                        page_sizes[str(pg_no)] = {
                            "width": getattr(sz, 'width', None),
                            "height": getattr(sz, 'height', None),
                        }

        # PPTX / DOCX: export markdown for structured rendering in the UI.
        markdown_text = None
        if ext in (".pptx", ".docx"):
            try:
                markdown_text = result.document.export_to_markdown()
            except Exception:
                logger.debug("export_to_markdown failed for %s — skipping", filename)

        logger.info(
            "Extraction completed: filename=%s pages=%d text_length=%d chunks=%d detected_language=%s title=%s author=%s",
            filename, final_pages, len(full_text), len(chunks_data), detected_lang,
            repr(title[:50]) if title else None,
            repr(author[:50]) if author else None,
        )

        return {
            "text": full_text,
            "page_count": final_pages,
            "title": title,
            "author": author,
            "detected_language": detected_lang,
            "chunks": chunks_data,
            "bbox_annotations": bbox_annotations,
            "page_sizes": page_sizes,
            "markdown_text": markdown_text,
        }
    except Exception as e:
        logger.error("Processing failed: %s", str(e), exc_info=True)
        return {"error": "Failed to process document. Please check the file and try again."}
    finally:
        if os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except OSError as e:
                logger.warning(f"Failed to delete temp file: {e}")

async def parse_document_async(file_bytes: bytes, filename: str) -> dict:
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(None, _parse_document_sync, file_bytes, filename)


async def warm_up_docling() -> None:
    """Pre-load Docling model weights at startup so the first upload is not slow."""
    loop = asyncio.get_running_loop()
    # Guard on → lightweight parsers only, docling is never used at runtime, so
    # don't load its weights (saves memory on the constrained host).
    if _ocr_blocked():
        logger.info("Docling warm-up skipped — BLOCK_OCR_DOCUMENTS on (lightweight parsers only)")
        return
    logger.info("Docling warm-up started — loading model weights")
    await loop.run_in_executor(None, lambda: _get_docling_converter(use_ocr=True))
    await loop.run_in_executor(None, _get_hybrid_chunker)
    logger.info("Docling warm-up complete")